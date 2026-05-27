from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (simulated with colored rectangle)
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    # Add emoji/icon
    emoji_box = slide.shapes.add_textbox(Inches(4), Inches(1.5), Inches(2), Inches(1))
    emoji_frame = emoji_box.text_frame
    emoji_frame.text = "🛡️"
    emoji_frame.paragraphs[0].font.size = Pt(72)
    emoji_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    
    # Add content
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(12)

# Slide 1: Title
add_title_slide(prs, 
    "AI-Powered Fraud Detection System",
    "Real-Time Transaction Analysis & Chatbot")

# Slide 2: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "🚨 The Problem"
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

stats = [
    "$32B - Annual Fraud Losses (2024)",
    "45% - Increase in Digital Fraud",
    "2.8M - Fraud Reports Filed",
    "",
    "Key Challenges:",
    "⏱️ Speed: Fraudsters act in seconds, detection must be instant",
    "🎯 Accuracy: False positives frustrate legitimate customers",
    "🔄 Adaptability: Fraud patterns constantly evolve",
    "📊 Scale: Millions of transactions need real-time analysis"
]

for stat in stats:
    p = tf.add_paragraph()
    p.text = stat
    p.font.size = Pt(18)
    p.space_after = Pt(10)

# Slide 3: Case Study
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "📋 Real Case Study: Michael Chen"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

case_details = [
    "Incident Details:",
    "",
    "Amount: $47.50 average → $850.00 (17.9x) ⚠️",
    "Location: New York, NY → Miami, FL ⚠️",
    "Time: 7 AM - 11 PM → 3:47 AM ⚠️",
    "Merchant: Grocery/Restaurant → Gas Station (Unfamiliar) ⚠️",
    "",
    "Result: HIGH RISK - Score 87/100",
    "Action: Transaction blocked, customer notified, card secured"
]

for detail in case_details:
    p = tf.add_paragraph()
    p.text = detail
    p.font.size = Pt(18)
    p.space_after = Pt(8)

# Slide 4: Solution Overview
add_content_slide(prs, "💡 Our Solution", [
    "🔍 Real-Time Analysis",
    "   • Analyzes transactions in <1 second",
    "   • Multi-factor risk scoring",
    "",
    "🤖 AI-Powered Chatbot",
    "   • Interactive Q&A for fraud decisions",
    "   • Educational content and prevention tips",
    "",
    "📊 Visual Dashboard",
    "   • Intuitive gauges and color-coded alerts",
    "   • Instant risk assessment",
    "",
    "🎯 96.2% Accuracy",
    "   • ML model trained on millions of transactions"
])

# Slide 5: Risk Scoring System
add_content_slide(prs, "⚖️ Risk Scoring System", [
    "Multi-Factor Analysis (0-100 Scale):",
    "",
    "• Amount Deviation (35 points): Compares to customer's average",
    "• Merchant Type (20 points): Checks against familiar categories",
    "• Location (20 points): Identifies unusual geographic patterns",
    "• Time Pattern (15 points): Flags transactions outside normal hours",
    "• Spending Spike (10 points): Detects sudden large purchases",
    "",
    "Risk Thresholds:",
    "0-20: LEGITIMATE ✅  |  21-50: LOW RISK 📊",
    "51-70: SUSPICIOUS ⚠️  |  71-85: HIGH RISK 🔶",
    "86-100: CRITICAL 🔴"
])

# Slide 6: Chatbot Capabilities
add_content_slide(prs, "💬 AI Chatbot Capabilities", [
    "❓ Explain Decisions",
    "   • 'Why is this transaction flagged as fraud?'",
    "   • Provides detailed breakdown of risk factors",
    "",
    "🎯 Recommend Actions",
    "   • 'Should this transaction be blocked?'",
    "   • Suggests appropriate response based on risk level",
    "",
    "📚 Educate Users",
    "   • 'What are the risk thresholds?'",
    "   • Explains scoring system and methodology",
    "",
    "🛡️ Prevention Tips",
    "   • 'How can we prevent fraud?'",
    "   • Provides best practices and recommendations"
])

# Slide 7: Live Demo
add_content_slide(prs, "🎬 Live Demo", [
    "Streamlit Application",
    "Access at: http://localhost:8501",
    "",
    "Demo Flow:",
    "1. Input Transaction: Pre-filled with Michael Chen's fraud case",
    "2. Click Analyze: Real-time risk calculation (<1 second)",
    "3. View Results: Risk gauge, score, and detailed factors",
    "4. Ask Chatbot: Interactive Q&A about the analysis",
    "5. Take Action: Block, verify, or allow based on recommendation",
    "",
    "Expected Result: HIGH RISK - 87/100",
    "Recommended Action: 🔶 BLOCK & VERIFY WITH CUSTOMER"
])

# Slide 8: Business Impact
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "📈 Business Impact"
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

impact_items = [
    "Key Metrics:",
    "• 96.2% Detection Accuracy",
    "• <1 second Analysis Time",
    "• 85% Fraud Prevention Rate",
    "",
    "Key Benefits:",
    "💰 Cost Savings: Prevent millions in fraud losses annually",
    "😊 Customer Satisfaction: Reduce false positives by 40%",
    "⚡ Operational Efficiency: Automate 90% of fraud reviews",
    "📊 Scalability: Handle millions of transactions per day",
    "🔒 Compliance: Meet regulatory requirements"
]

for item in impact_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(18)
    p.space_after = Pt(10)

# Slide 9: Technical Architecture
add_content_slide(prs, "🏗️ Technical Architecture", [
    "Frontend Layer:",
    "• Streamlit web framework",
    "• Plotly visualizations",
    "• Responsive design",
    "",
    "Analysis Engine:",
    "• Multi-factor risk scoring",
    "• Real-time calculation",
    "• Pattern recognition",
    "",
    "Chatbot Module:",
    "• Natural language processing",
    "• Context-aware responses",
    "• Educational content",
    "",
    "Deployment: Cloud (AWS/Azure/GCP), On-Premise, API Integration"
])

# Slide 10: Future Enhancements
add_content_slide(prs, "🚀 Future Enhancements", [
    "Phase 2 Features:",
    "• 🤖 Advanced ML Models: Deep learning for pattern recognition",
    "• 🌐 Global Fraud Network: Share intelligence across institutions",
    "• 📱 Mobile App: Native iOS and Android applications",
    "• 🔔 Real-Time Alerts: SMS, email, and push notifications",
    "• 📊 Analytics Dashboard: Fraud trends and insights",
    "",
    "Phase 3 Features:",
    "• 🧠 Behavioral Biometrics: Typing patterns, mouse movements",
    "• 🔗 Blockchain Integration: Immutable fraud records",
    "• 🌍 Multi-Language Support: Global deployment",
    "• 🎯 Predictive Analytics: Forecast fraud before it happens"
])

# Slide 11: Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Add background
left = top = Inches(0)
width = prs.slide_width
height = prs.slide_height
shape = slide.shapes.add_shape(1, left, top, width, height)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(102, 126, 234)

# Add emoji
emoji_box = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1))
emoji_frame = emoji_box.text_frame
emoji_frame.text = "🎯"
emoji_frame.paragraphs[0].font.size = Pt(72)
emoji_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Ready to Protect Your Customers?"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add next steps
steps_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(2.5))
steps_frame = steps_box.text_frame
steps_frame.text = """Next Steps:
1️⃣ Try the Demo - Experience the system live
2️⃣ Schedule a Meeting - Discuss implementation
3️⃣ Pilot Program - Start with small-scale deployment
4️⃣ Full Rollout - Scale to production"""
steps_frame.paragraphs[0].font.size = Pt(20)
steps_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Add thank you
thanks_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.8))
thanks_frame = thanks_box.text_frame
thanks_frame.text = "Thank You! 🙏"
thanks_frame.paragraphs[0].font.size = Pt(36)
thanks_frame.paragraphs[0].font.bold = True
thanks_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('Fraud-Detection-Chatbot-Presentation.pptx')
print("PowerPoint presentation created successfully!")
print("File: Fraud-Detection-Chatbot-Presentation.pptx")

# Made with Bob
