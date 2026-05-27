from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_flow_diagram_slide(prs, title):
    """Add a blank slide for flow diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_shape_with_text(slide, shape_type, left, top, width, height, text, bg_color, text_color=RGBColor(255, 255, 255)):
    """Add a shape with text"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.5)
    
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = text_color
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = 1  # Middle
    
    return shape

def add_arrow(slide, x1, y1, x2, y2, label=""):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight connector
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(2)
    
    if label:
        # Add label near the arrow
        label_box = slide.shapes.add_textbox(
            Inches((x1.inches + x2.inches) / 2 - 0.5),
            Inches((y1.inches + y2.inches) / 2 - 0.2),
            Inches(1),
            Inches(0.3)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].font.size = Pt(10)
        label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_chatbot_flow_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "🤖 Fraud Detection Chatbot",
        "Step-by-Step Flow Diagram"
    )
    
    # Slide 2: Overall Architecture
    slide2 = add_flow_diagram_slide(prs, "Overall System Architecture")
    
    # User
    add_shape_with_text(
        slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(2), Inches(0.8),
        "👤 User", RGBColor(52, 152, 219)
    )
    
    # Web Interface
    add_shape_with_text(
        slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(1.5), Inches(3), Inches(0.8),
        "🌐 Streamlit Web Interface", RGBColor(46, 204, 113)
    )
    
    # Fraud Detection Engine
    add_shape_with_text(
        slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(3), Inches(3), Inches(0.8),
        "🛡️ Fraud Detection Engine", RGBColor(231, 76, 60)
    )
    
    # Chatbot AI
    add_shape_with_text(
        slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(4.5), Inches(3), Inches(0.8),
        "💬 AI Chatbot", RGBColor(155, 89, 182)
    )
    
    # Database
    add_shape_with_text(
        slide2, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(3), Inches(2), Inches(1.5),
        "📊 Transaction\nData", RGBColor(241, 196, 15)
    )
    
    # Arrows
    add_arrow(slide2, Inches(2.5), Inches(1.9), Inches(3.5), Inches(1.9))
    add_arrow(slide2, Inches(5), Inches(2.3), Inches(5), Inches(3))
    add_arrow(slide2, Inches(5), Inches(3.8), Inches(5), Inches(4.5))
    add_arrow(slide2, Inches(6.5), Inches(3.4), Inches(7.5), Inches(3.4))
    
    # Slide 3: User Interaction Flow
    slide3 = add_flow_diagram_slide(prs, "Step 1: User Interaction Flow")
    
    y_pos = 1.5
    step_height = 0.7
    step_gap = 0.3
    
    # Step 1
    add_shape_with_text(
        slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(step_height),
        "1️⃣ User Opens Fraud Detection App", RGBColor(52, 152, 219)
    )
    add_arrow(slide3, Inches(5), Inches(y_pos + step_height), Inches(5), Inches(y_pos + step_height + step_gap))
    
    y_pos += step_height + step_gap
    
    # Step 2
    add_shape_with_text(
        slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(step_height),
        "2️⃣ User Enters Transaction Details", RGBColor(46, 204, 113)
    )
    add_arrow(slide3, Inches(5), Inches(y_pos + step_height), Inches(5), Inches(y_pos + step_height + step_gap))
    
    y_pos += step_height + step_gap
    
    # Step 3
    add_shape_with_text(
        slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(step_height),
        "3️⃣ User Clicks 'Analyze Fraud Risk'", RGBColor(155, 89, 182)
    )
    add_arrow(slide3, Inches(5), Inches(y_pos + step_height), Inches(5), Inches(y_pos + step_height + step_gap))
    
    y_pos += step_height + step_gap
    
    # Step 4
    add_shape_with_text(
        slide3, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(step_height),
        "4️⃣ System Displays Risk Analysis", RGBColor(231, 76, 60)
    )
    
    # Slide 4: Fraud Detection Process
    slide4 = add_flow_diagram_slide(prs, "Step 2: Fraud Detection Analysis Process")
    
    y_pos = 1.5
    
    # Input
    add_shape_with_text(
        slide4, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(0.6),
        "📥 Transaction Input Received", RGBColor(52, 152, 219)
    )
    add_arrow(slide4, Inches(5), Inches(y_pos + 0.6), Inches(5), Inches(y_pos + 0.9))
    
    y_pos += 0.9
    
    # Analysis steps
    analyses = [
        ("💰 Amount Deviation Analysis", RGBColor(46, 204, 113)),
        ("🏪 Merchant Type Analysis", RGBColor(155, 89, 182)),
        ("📍 Location Analysis", RGBColor(241, 196, 15)),
        ("🕐 Time Pattern Analysis", RGBColor(230, 126, 34)),
        ("📊 Velocity Check", RGBColor(231, 76, 60))
    ]
    
    for i, (text, color) in enumerate(analyses):
        add_shape_with_text(
            slide4, MSO_SHAPE.RECTANGLE,
            Inches(2), Inches(y_pos), Inches(6), Inches(0.5),
            text, color
        )
        if i < len(analyses) - 1:
            add_arrow(slide4, Inches(5), Inches(y_pos + 0.5), Inches(5), Inches(y_pos + 0.7))
        y_pos += 0.7
    
    # Output
    add_arrow(slide4, Inches(5), Inches(y_pos - 0.2), Inches(5), Inches(y_pos + 0.1))
    add_shape_with_text(
        slide4, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos + 0.1), Inches(6), Inches(0.6),
        "📤 Risk Score Generated (0-100)", RGBColor(192, 57, 43)
    )
    
    # Slide 5: Risk Scoring Logic
    slide5 = add_flow_diagram_slide(prs, "Step 3: Risk Scoring Logic")
    
    # Decision tree
    add_shape_with_text(
        slide5, MSO_SHAPE.DIAMOND,
        Inches(3.5), Inches(1.5), Inches(3), Inches(0.8),
        "Risk Score?", RGBColor(52, 152, 219)
    )
    
    # Branches
    risk_levels = [
        (1, "0-20", "✅ LEGITIMATE", RGBColor(46, 204, 113)),
        (2.5, "21-50", "📊 LOW RISK", RGBColor(52, 152, 219)),
        (4, "51-70", "⚠️ SUSPICIOUS", RGBColor(241, 196, 15)),
        (5.5, "71-85", "🔶 HIGH RISK", RGBColor(230, 126, 34)),
        (7, "86-100", "🔴 CRITICAL", RGBColor(231, 76, 60))
    ]
    
    for i, (x_pos, score_range, label, color) in enumerate(risk_levels):
        add_shape_with_text(
            slide5, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(3.5), Inches(1.3), Inches(0.7),
            f"{score_range}\n{label}", color
        )
        # Arrow from decision to outcome
        add_arrow(slide5, Inches(5), Inches(2.3), Inches(x_pos + 0.65), Inches(3.5))
    
    # Slide 6: Chatbot Interaction Flow
    slide6 = add_flow_diagram_slide(prs, "Step 4: Chatbot Interaction Flow")
    
    y_pos = 1.5
    
    # User asks question
    add_shape_with_text(
        slide6, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(0.7),
        "💬 User Asks Question", RGBColor(52, 152, 219)
    )
    add_arrow(slide6, Inches(5), Inches(y_pos + 0.7), Inches(5), Inches(y_pos + 1))
    
    y_pos += 1
    
    # Question analysis
    add_shape_with_text(
        slide6, MSO_SHAPE.HEXAGON,
        Inches(2), Inches(y_pos), Inches(6), Inches(0.7),
        "🔍 Analyze Question Intent", RGBColor(155, 89, 182)
    )
    add_arrow(slide6, Inches(5), Inches(y_pos + 0.7), Inches(5), Inches(y_pos + 1))
    
    y_pos += 1
    
    # Question types
    question_types = [
        ("Why flagged?", Inches(1), RGBColor(46, 204, 113)),
        ("Should block?", Inches(3.5), RGBColor(241, 196, 15)),
        ("Prevention tips?", Inches(6), RGBColor(230, 126, 34))
    ]
    
    for text, x_pos, color in question_types:
        add_shape_with_text(
            slide6, MSO_SHAPE.RECTANGLE,
            x_pos, Inches(y_pos), Inches(2), Inches(0.6),
            text, color
        )
        add_arrow(slide6, Inches(5), Inches(y_pos - 0.3), x_pos + Inches(1), Inches(y_pos))
    
    y_pos += 0.9
    
    # Generate response
    add_shape_with_text(
        slide6, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(0.7),
        "🤖 Generate Contextual Response", RGBColor(231, 76, 60)
    )
    add_arrow(slide6, Inches(5), Inches(y_pos + 0.7), Inches(5), Inches(y_pos + 1))
    
    y_pos += 1
    
    # Display response
    add_shape_with_text(
        slide6, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(y_pos), Inches(6), Inches(0.7),
        "💬 Display Response to User", RGBColor(52, 152, 219)
    )
    
    # Slide 7: Response Generation Logic
    slide7 = add_flow_diagram_slide(prs, "Step 5: Chatbot Response Generation")
    
    # Input
    add_shape_with_text(
        slide7, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(1.5), Inches(4), Inches(0.6),
        "📥 User Question Received", RGBColor(52, 152, 219)
    )
    add_arrow(slide7, Inches(5), Inches(2.1), Inches(5), Inches(2.5))
    
    # Decision
    add_shape_with_text(
        slide7, MSO_SHAPE.DIAMOND,
        Inches(3.5), Inches(2.5), Inches(3), Inches(0.8),
        "Question Type?", RGBColor(155, 89, 182)
    )
    
    # Response types
    responses = [
        (0.5, "Explain Risk", "Show risk factors\n& score breakdown", RGBColor(46, 204, 113)),
        (2.5, "Action Advice", "Recommend\nblock/allow/verify", RGBColor(241, 196, 15)),
        (4.5, "Prevention", "Provide fraud\nprevention tips", RGBColor(230, 126, 34)),
        (6.5, "Thresholds", "Explain risk\nscore ranges", RGBColor(231, 76, 60))
    ]
    
    for x_pos, title, desc, color in responses:
        add_shape_with_text(
            slide7, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(4), Inches(1.8), Inches(1),
            f"{title}\n\n{desc}", color
        )
        add_arrow(slide7, Inches(5), Inches(3.3), Inches(x_pos + 0.9), Inches(4))
    
    # Slide 8: Complete User Journey
    slide8 = add_flow_diagram_slide(prs, "Complete User Journey")
    
    journey_steps = [
        ("🚀 Start", "User opens app", RGBColor(52, 152, 219)),
        ("📝 Input", "Enter transaction", RGBColor(46, 204, 113)),
        ("🔍 Analyze", "Click analyze", RGBColor(155, 89, 182)),
        ("📊 Results", "View risk score", RGBColor(241, 196, 15)),
        ("💬 Ask", "Ask chatbot", RGBColor(230, 126, 34)),
        ("💡 Learn", "Get insights", RGBColor(231, 76, 60)),
        ("✅ Decide", "Take action", RGBColor(46, 204, 113))
    ]
    
    x_start = 1
    y_pos = 3
    
    for i, (icon, text, color) in enumerate(journey_steps):
        add_shape_with_text(
            slide8, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_start + i * 1.2), Inches(y_pos), Inches(1), Inches(1),
            f"{icon}\n\n{text}", color
        )
        if i < len(journey_steps) - 1:
            add_arrow(
                slide8,
                Inches(x_start + i * 1.2 + 1),
                Inches(y_pos + 0.5),
                Inches(x_start + (i + 1) * 1.2),
                Inches(y_pos + 0.5)
            )
    
    # Slide 9: Key Features
    slide9 = add_flow_diagram_slide(prs, "Key Features Summary")
    
    features = [
        ("⚡ Real-time Analysis", "< 1 second response", Inches(1), Inches(2)),
        ("🎯 96.2% Accuracy", "ML-powered detection", Inches(4), Inches(2)),
        ("💬 Smart Chatbot", "Context-aware responses", Inches(7), Inches(2)),
        ("📊 Visual Dashboard", "Interactive gauges", Inches(1), Inches(4)),
        ("🔔 Instant Alerts", "Multi-channel notify", Inches(4), Inches(4)),
        ("🛡️ Multi-factor", "5 risk indicators", Inches(7), Inches(4))
    ]
    
    for title, desc, x, y in features:
        add_shape_with_text(
            slide9, MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(2.5), Inches(1.2),
            f"{title}\n\n{desc}", RGBColor(52, 152, 219)
        )
    
    # Slide 10: Summary
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Summary"
    
    content = slide10.placeholders[1].text_frame
    content.text = "Fraud Detection Chatbot Flow:\n\n"
    
    summary_points = [
        "1. User enters transaction details in web interface",
        "2. System analyzes 5 risk factors in real-time",
        "3. Risk score (0-100) calculated and displayed",
        "4. Visual dashboard shows results with gauges",
        "5. User can ask chatbot questions about analysis",
        "6. Chatbot provides context-aware responses",
        "7. System recommends action based on risk level",
        "8. Complete audit trail maintained"
    ]
    
    for point in summary_points:
        p = content.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
    
    # Save presentation
    prs.save('Fraud-Detection-Chatbot-Flow-Diagram.pptx')
    print("SUCCESS: PowerPoint presentation created: Fraud-Detection-Chatbot-Flow-Diagram.pptx")

if __name__ == "__main__":
    create_chatbot_flow_presentation()

# Made with Bob
