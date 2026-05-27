from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(118, 75, 162)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 1: Title
add_title_slide(prs, 
    "AI-Powered Fraud Detection System",
    "Real-Time Transaction Analysis & Chatbot")

# Slide 2: Fraud Statistics with Bar Chart
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
title = slide.shapes.title
title.text = "Fraud Statistics - The Growing Problem"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create bar chart
chart_data = CategoryChartData()
chart_data.categories = ['2020', '2021', '2022', '2023', '2024']
chart_data.add_series('Fraud Losses ($B)', (18, 22, 26, 29, 32))

x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# Slide 3: Fraud Types Distribution - Pie Chart
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Fraud Types Distribution (2024)"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create pie chart
chart_data = CategoryChartData()
chart_data.categories = ['Card Fraud', 'Identity Theft', 'Account Takeover', 'Wire Transfer', 'Other']
chart_data.add_series('Percentage', (35, 25, 20, 12, 8))

x, y, cx, cy = Inches(2), Inches(1.5), Inches(6), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT

# Slide 4: Michael Chen Case - Transaction Comparison
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Case Study: Michael Chen - Transaction Analysis"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create comparison chart
chart_data = CategoryChartData()
chart_data.categories = ['Amount', 'Time (Hour)', 'Distance (miles)']
chart_data.add_series('Normal Pattern', (47.5, 15, 5))
chart_data.add_series('Fraudulent Transaction', (850, 3.8, 1200))

x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM

# Slide 5: Risk Score Distribution
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Risk Score Distribution - Last 1000 Transactions"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create area chart
chart_data = CategoryChartData()
chart_data.categories = ['0-20', '21-50', '51-70', '71-85', '86-100']
chart_data.add_series('Transactions', (750, 180, 45, 20, 5))

x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data
).chart

chart.has_legend = False

# Slide 6: Detection Accuracy Over Time
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Model Accuracy Improvement Over Time"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create line chart
chart_data = CategoryChartData()
chart_data.categories = ['Q1 2023', 'Q2 2023', 'Q3 2023', 'Q4 2023', 'Q1 2024', 'Q2 2024']
chart_data.add_series('Accuracy %', (88.5, 90.2, 92.8, 94.1, 95.3, 96.2))

x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

chart.has_legend = False

# Slide 7: Risk Factor Contribution
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Risk Factor Contribution (Max Points)"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create horizontal bar chart
chart_data = CategoryChartData()
chart_data.categories = ['Amount Deviation', 'Merchant Type', 'Location', 'Time Pattern', 'Spending Spike']
chart_data.add_series('Points', (35, 20, 20, 15, 10))

x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
).chart

chart.has_legend = False

# Slide 8: Monthly Fraud Prevention Impact
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Monthly Fraud Prevention Impact"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create stacked column chart
chart_data = CategoryChartData()
chart_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
chart_data.add_series('Prevented ($M)', (2.5, 3.1, 2.8, 3.5, 4.2, 3.9))
chart_data.add_series('Detected ($M)', (0.5, 0.4, 0.6, 0.3, 0.4, 0.5))

x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM

# Slide 9: Response Time Analysis
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "System Response Time Distribution"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create chart
chart_data = CategoryChartData()
chart_data.categories = ['<0.5s', '0.5-1s', '1-2s', '2-3s', '>3s']
chart_data.add_series('Transactions (%)', (85, 12, 2, 0.8, 0.2))

x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

chart.has_legend = False

# Slide 10: ROI Analysis
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Return on Investment (ROI) - 6 Month Period"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

# Create chart
chart_data = CategoryChartData()
chart_data.categories = ['Implementation Cost', 'Fraud Prevented', 'Net Savings']
chart_data.add_series('Amount ($M)', (2.5, 20.5, 18.0))

x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

chart.has_legend = False

# Slide 11: Key Metrics Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Key Performance Metrics"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add metric boxes
metrics = [
    ("96.2%", "Detection Accuracy", 1, 1.2),
    ("<1 sec", "Analysis Time", 4, 1.2),
    ("85%", "Prevention Rate", 7, 1.2),
    ("$20.5M", "Fraud Prevented", 1, 3.5),
    ("720:1", "ROI Ratio", 4, 3.5),
    ("99.9%", "Uptime", 7, 3.5),
    ("2.8%", "False Positive Rate", 1, 5.8),
    ("1M+", "Daily Transactions", 4, 5.8),
    ("24/7", "Monitoring", 7, 5.8)
]

for value, label, x, y in metrics:
    # Value box
    value_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.5), Inches(0.8))
    value_frame = value_box.text_frame
    value_frame.text = value
    value_frame.paragraphs[0].font.size = Pt(32)
    value_frame.paragraphs[0].font.bold = True
    value_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Label box
    label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.8), Inches(2.5), Inches(0.4))
    label_frame = label_box.text_frame
    label_frame.text = label
    label_frame.paragraphs[0].font.size = Pt(14)
    label_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 12: Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Ready to Protect Your Customers?"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

content_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(2.5))
content_frame = content_box.text_frame
content_frame.text = """Next Steps:
1. Try the Live Demo at localhost:8501
2. Schedule Implementation Meeting
3. Start Pilot Program
4. Full Production Rollout

Thank You!"""
content_frame.paragraphs[0].font.size = Pt(20)
content_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
content_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('Fraud-Detection-Presentation-With-Graphs.pptx')
print("PowerPoint presentation with graphs created successfully!")
print("File: Fraud-Detection-Presentation-With-Graphs.pptx")

# Made with Bob
